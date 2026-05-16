#!/usr/bin/env python3
"""Build a 1280x720 image-per-slide PPTX from /tmp/deck-pngs/slide-*.png."""
import glob, os
from pptx import Presentation
from pptx.util import Emu

PNG_DIR = "/tmp/deck-pngs"
OUT = "investor-materials/GUBER-Investor-Deck.pptx"

prs = Presentation()
prs.slide_width = Emu(12192000)   # 13.333 in (16:9)
prs.slide_height = Emu(6858000)   # 7.5 in
blank = prs.slide_layouts[6]

for png in sorted(glob.glob(os.path.join(PNG_DIR, "slide-*.png"))):
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(png, 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(OUT)
size_mb = os.path.getsize(OUT) / (1024*1024)
print(f"wrote {OUT} with {len(prs.slides)} slides ({size_mb:.1f} MB)")
